"""
File Splitter Utility for handling large file splitting with semantic integrity.
"""

import json
import os
import uuid
import warnings
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import defusedxml.ElementTree as ET
import markdown
import pandas as pd
import pypdf
import tiktoken
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation

from app.be_core.config import settings
from app.be_core.logger import logger

# Suppress pandas DtypeWarning for mixed types in CSV files
warnings.filterwarnings(
    "ignore", message=".*mixed types.*", category=pd.errors.DtypeWarning
)

try:
    import chardet

    CHARDET_AVAILABLE = True
except ImportError:
    CHARDET_AVAILABLE = False

try:
    import fitz  # PyMuPDF

    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    import pdfplumber

    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False


# Define sample sizes and thresholds for different file sizes
SMALL_FILE_THRESHOLD = 1 * 1024 * 1024  # 1MB
MEDIUM_FILE_THRESHOLD = 10 * 1024 * 1024  # 10MB
LARGE_FILE_THRESHOLD = 50 * 1024 * 1024  # 50MB
SAMPLE_SIZE_SMALL = 20 * 1024  # 20KB for small files
SAMPLE_SIZE_MEDIUM = 50 * 1024  # 50KB for medium files
SAMPLE_SIZE_LARGE = 100 * 1024  # 100KB for large files
MAX_SAMPLES = 100  # Number of samples to take from different parts of the file


class PDFCorruptionError(Exception):
    """Exception raised when a PDF is detected as corrupted or damaged."""

    def __init__(self, message: str, corruption_details: str = ""):
        self.message = message
        self.corruption_details = corruption_details
        super().__init__(self.message)


class FileSplitter:
    """Utility for splitting large files into smaller chunks while preserving semantic meaning."""

    def __init__(
        self,
        max_tokens_per_split: Optional[int] = None,
        min_tokens_per_split: Optional[int] = None,
    ):
        """
        Initialize the FileSplitter with configuration.

        Args:
            max_tokens_per_split: Maximum tokens allowed per split
            min_tokens_per_split: Minimum tokens allowed per split to ensure meaningful chunks
        """
        self.max_tokens_per_split = (
            max_tokens_per_split or settings.MAX_TOKENS_PER_SPLIT
        )
        self.min_tokens_per_split = min_tokens_per_split or settings.MIN_SPLIT_SIZE

        logger.info(
            f"Initializing FileSplitter with max_tokens_per_split={self.max_tokens_per_split}, min_tokens_per_split={self.min_tokens_per_split}"
        )

        # Add file size thresholds as class attributes
        self.SMALL_FILE_THRESHOLD = SMALL_FILE_THRESHOLD
        self.MEDIUM_FILE_THRESHOLD = MEDIUM_FILE_THRESHOLD
        self.LARGE_FILE_THRESHOLD = LARGE_FILE_THRESHOLD
        self.SAMPLE_SIZE_SMALL = SAMPLE_SIZE_SMALL
        self.SAMPLE_SIZE_MEDIUM = SAMPLE_SIZE_MEDIUM
        self.SAMPLE_SIZE_LARGE = SAMPLE_SIZE_LARGE
        self.MAX_SAMPLES = MAX_SAMPLES

        # Ensure the temp splits directory exists
        os.makedirs(settings.TEMP_SPLITS_DIR, exist_ok=True)
        logger.debug(f"Temporary splits directory: {settings.TEMP_SPLITS_DIR}")

        self.encoder = tiktoken.get_encoding("cl100k_base")  # GPT-4 encoding
        logger.debug("Initialized tiktoken encoder with cl100k_base encoding")

    def count_tokens(self, text: str) -> int:
        """Count the number of tokens in a text using tiktoken."""
        token_count = len(self.encoder.encode(text))
        # logger.debug(f"Counted {token_count} tokens in text of length {len(text)}")
        return token_count

    def _validate_pdf_integrity(self, file_path: Path) -> None:
        """
        Comprehensive PDF corruption detection using multiple validation techniques.

        Raises PDFCorruptionError if the PDF is corrupted and should not be processed.

        Args:
            file_path: Path to the PDF file to validate

        Raises:
            PDFCorruptionError: If PDF is corrupted or damaged
            FileNotFoundError: If file doesn't exist
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PDF file not found: {file_path}")

        logger.info(f"Validating PDF integrity for: {file_path}")

        # Track corruption indicators
        corruption_issues = []

        # 1. Basic file structure validation
        try:
            self._validate_pdf_file_structure(file_path, corruption_issues)
        except Exception as e:
            corruption_issues.append(f"File structure validation failed: {str(e)}")

        # 2. Try validation with each available PDF library
        library_results = self._validate_with_multiple_libraries(
            file_path, corruption_issues
        )

        # 3. Analyze results and decide if PDF should be rejected
        self._analyze_corruption_results(file_path, corruption_issues, library_results)

    def _validate_pdf_file_structure(
        self, file_path: Path, corruption_issues: List[str]
    ) -> None:
        """Validate basic PDF file structure."""
        file_size = file_path.stat().st_size

        # Check file size (PDF should be at least a few KB)
        if file_size < 1024:  # Less than 1KB is suspicious
            corruption_issues.append(
                f"File too small ({file_size} bytes) - likely corrupted"
            )

        # Check PDF header and trailer
        try:
            with file_path.open("rb") as f:
                # Check PDF header
                header = f.read(8)
                if not header.startswith(b"%PDF-"):
                    corruption_issues.append(
                        "Invalid PDF header - missing '%PDF-' signature"
                    )

                # Check for basic PDF structure markers
                f.seek(0)
                content = f.read(
                    min(8192, file_size)
                )  # Read first 8KB or whole file if smaller

                # Look for essential PDF elements
                if b"%%EOF" not in content and file_size > 8192:
                    # For larger files, check the end
                    f.seek(max(0, file_size - 1024))
                    trailer_content = f.read()
                    if b"%%EOF" not in trailer_content:
                        corruption_issues.append(
                            "Missing PDF end-of-file marker (%%EOF)"
                        )

                # Check for xref table or stream
                f.seek(0)
                full_content = f.read()
                if b"xref" not in full_content and b"/XRefStm" not in full_content:
                    corruption_issues.append("Missing PDF cross-reference table")

        except Exception as e:
            corruption_issues.append(f"File structure analysis failed: {str(e)}")

    def _validate_with_multiple_libraries(
        self, file_path: Path, corruption_issues: List[str]
    ) -> Dict[str, Any]:
        """Validate PDF with multiple libraries to detect corruption patterns."""
        library_results = {
            "pypdf": {"success": False, "pages": 0, "errors": []},
            "pymupdf": {"success": False, "pages": 0, "errors": []},
            "pdfplumber": {"success": False, "pages": 0, "errors": []},
        }

        # Test pypdf
        try:
            with open(str(file_path), "rb") as file:
                reader = pypdf.PdfReader(file)
                page_count = len(reader.pages)

                # Try to access first and last pages to check for corruption
                if page_count > 0:
                    try:
                        first_page = reader.pages[0]
                        # Try to extract basic information
                        first_page.extract_text()
                        if page_count > 1:
                            last_page = reader.pages[-1]
                            last_page.extract_text()

                        library_results["pypdf"] = {
                            "success": True,
                            "pages": page_count,
                            "errors": [],
                        }

                    except Exception as e:
                        library_results["pypdf"]["errors"].append(
                            f"Page access failed: {str(e)}"
                        )

        except Exception as e:
            library_results["pypdf"]["errors"].append(
                f"Reader creation failed: {str(e)}"
            )

        # Test PyMuPDF if available
        if PYMUPDF_AVAILABLE:
            try:
                doc = fitz.open(str(file_path))
                page_count = doc.page_count

                if page_count > 0:
                    try:
                        # Test page access
                        first_page = doc.load_page(0)
                        first_page.get_text()
                        if page_count > 1:
                            last_page = doc.load_page(-1)
                            last_page.get_text()

                        library_results["pymupdf"] = {
                            "success": True,
                            "pages": page_count,
                            "errors": [],
                        }

                    except Exception as e:
                        library_results["pymupdf"]["errors"].append(
                            f"Page access failed: {str(e)}"
                        )
                    finally:
                        try:
                            doc.close()
                        except Exception as e:
                            logger.warning(
                                f"Failed to close PyMuPDF document: {str(e)}"
                            )

            except Exception as e:
                library_results["pymupdf"]["errors"].append(
                    f"Document opening failed: {str(e)}"
                )

        # Test pdfplumber if available
        if PDFPLUMBER_AVAILABLE:
            try:
                with pdfplumber.open(str(file_path)) as pdf:
                    page_count = len(pdf.pages)

                    if page_count > 0:
                        try:
                            # Test page access
                            first_page = pdf.pages[0]
                            first_page.extract_text()
                            if page_count > 1:
                                last_page = pdf.pages[-1]
                                last_page.extract_text()

                            library_results["pdfplumber"] = {
                                "success": True,
                                "pages": page_count,
                                "errors": [],
                            }

                        except Exception as e:
                            library_results["pdfplumber"]["errors"].append(
                                f"Page access failed: {str(e)}"
                            )

            except Exception as e:
                library_results["pdfplumber"]["errors"].append(
                    f"PDF opening failed: {str(e)}"
                )

        return library_results

    def _analyze_corruption_results(
        self,
        file_path: Path,
        corruption_issues: List[str],
        library_results: Dict[str, Any],
    ) -> None:
        """Analyze corruption detection results and raise exception if PDF should be rejected."""

        # Count successful library validations
        successful_libraries = [
            lib for lib, result in library_results.items() if result["success"]
        ]
        failed_libraries = [
            lib for lib, result in library_results.items() if not result["success"]
        ]

        # Check for page count inconsistencies (major red flag)
        page_counts = [
            result["pages"] for result in library_results.values() if result["success"]
        ]
        if len(set(page_counts)) > 1 and len(page_counts) > 1:
            corruption_issues.append(
                f"Page count inconsistency across libraries: {dict(zip(successful_libraries, page_counts))}"
            )

        # Severity assessment - be more strict about critical structural issues
        critical_keywords = [
            "header",
            "eof",
            "xref",
            "structure",
            "inconsistency",
            "signature",
        ]
        critical_issues = len(
            [
                issue
                for issue in corruption_issues
                if any(word in issue.lower() for word in critical_keywords)
            ]
        )

        # Check for specific critical issues that should always cause rejection
        blocking_issues = [
            issue
            for issue in corruption_issues
            if any(
                blocker in issue.lower()
                for blocker in [
                    "invalid pdf header",
                    "missing",
                    "signature",
                    "not a pdf",
                ]
            )
        ]

        # Decision logic - stricter rules for rejection
        total_issues = len(corruption_issues)

        # If no libraries can handle the PDF, it's definitely corrupted
        if not successful_libraries:
            all_errors = []
            for lib, result in library_results.items():
                if result["errors"]:
                    all_errors.extend([f"{lib}: {error}" for error in result["errors"]])

            error_details = "; ".join(corruption_issues + all_errors)
            raise PDFCorruptionError(
                "PDF is corrupted - no library can process it successfully",
                error_details,
            )

        # If we have any blocking issues (like invalid header), reject immediately
        if blocking_issues:
            error_details = "; ".join(corruption_issues)
            raise PDFCorruptionError(
                "PDF has critical blocking issues that prevent processing",
                error_details,
            )

        # If we have multiple critical structural issues, reject
        if critical_issues >= 2:
            error_details = "; ".join(corruption_issues)
            raise PDFCorruptionError(
                f"PDF has multiple critical structural corruption issues ({critical_issues} critical issues found)",
                error_details,
            )

        # If majority of libraries fail and we have structural issues, reject
        if len(failed_libraries) > len(successful_libraries) and total_issues >= 1:
            error_details = "; ".join(corruption_issues)
            raise PDFCorruptionError(
                f"PDF shows significant corruption signs - {len(failed_libraries)}/{len(library_results)} libraries failed",
                error_details,
            )

        # STRICTER: If we have page count inconsistencies (indicates structural corruption), reject
        has_page_inconsistency = any(
            "inconsistency" in issue.lower() for issue in corruption_issues
        )
        if has_page_inconsistency:
            error_details = "; ".join(corruption_issues)
            raise PDFCorruptionError(
                "PDF has structural corruption - page count inconsistencies indicate corrupted PDF structure",
                error_details,
            )

        # STRICTER: If 2+ libraries fail (even if 1 succeeds), reject as unreliable
        if len(failed_libraries) >= 2:
            error_details = "; ".join(
                corruption_issues
                + [
                    f"{lib}: {'; '.join(library_results[lib]['errors'])}"
                    for lib in failed_libraries
                ]
            )
            raise PDFCorruptionError(
                f"PDF has significant corruption - multiple libraries ({len(failed_libraries)}/{len(library_results)}) failed to process it reliably",
                error_details,
            )

        # If we get here, the PDF is likely processable despite minor issues
        if corruption_issues:
            logger.warning(
                f"PDF has minor issues but is processable: {'; '.join(corruption_issues)}"
            )
        else:
            logger.info(
                f"PDF validation successful - {len(successful_libraries)} libraries can process it"
            )

    def detect_file_encoding(self, file_path: Path, sample_size: int = 8192) -> str:
        """
        Detect the encoding of a file using chardet if available, with robust fallback strategies.

        Args:
            file_path: Path to the file
            sample_size: Number of bytes to sample for encoding detection

        Returns:
            Detected encoding string
        """
        logger.debug(f"Detecting encoding for file: {file_path}")

        try:
            # Read a sample of the file in binary mode
            with file_path.open("rb") as f:
                # Read up to sample_size bytes for detection
                sample_data = f.read(sample_size)

            if not sample_data:
                logger.warning(f"File {file_path} is empty, defaulting to UTF-8")
                return "utf-8"

            # Try chardet if available
            if CHARDET_AVAILABLE:
                result = chardet.detect(sample_data)
                detected_encoding = result.get("encoding", "utf-8")
                confidence = result.get("confidence", 0.0)

                logger.info(
                    f"Chardet detected encoding: {detected_encoding} (confidence: {confidence:.2f})"
                )

                # Be skeptical of ASCII detection - it's often wrong for large files
                if detected_encoding and detected_encoding.lower() == "ascii":
                    logger.warning(
                        "ASCII detected by chardet, but this is often incorrect for data files. Using fallback detection."
                    )
                    return self._detect_encoding_fallback(sample_data, file_path)

                # Use detected encoding if confidence is reasonable and it's not ASCII
                if (
                    confidence >= 0.7
                    and detected_encoding
                    and detected_encoding.lower() != "ascii"
                ):
                    # Validate the detected encoding by trying to decode a larger sample
                    if self._validate_encoding(file_path, detected_encoding):
                        return detected_encoding.lower()
                    else:
                        logger.warning(
                            f"Detected encoding {detected_encoding} failed validation, using fallbacks"
                        )

                logger.warning(
                    f"Low confidence ({confidence:.2f}) or failed validation for detected encoding {detected_encoding}, will try fallbacks"
                )
            else:
                logger.info("Chardet not available, using fallback detection")

            # Fallback: try common encodings
            return self._detect_encoding_fallback(sample_data, file_path)

        except Exception as e:
            logger.error(f"Error detecting encoding for {file_path}: {str(e)}")
            # Final fallback
            return "utf-8"

    def _validate_encoding(
        self, file_path: Path, encoding: str, validation_size: int = 32768
    ) -> bool:
        """
        Validate an encoding by trying to read a larger sample of the file.

        Args:
            file_path: Path to the file
            encoding: Encoding to validate
            validation_size: Size of sample to read for validation

        Returns:
            True if encoding can decode the sample without errors
        """
        try:
            with file_path.open("r", encoding=encoding, errors="strict") as f:
                # Try to read a larger sample to validate the encoding
                f.read(validation_size)
            return True
        except (UnicodeDecodeError, UnicodeError, LookupError):
            return False

    def _detect_encoding_fallback(self, sample_data: bytes, file_path: Path) -> str:
        """
        Fallback encoding detection by trying common encodings.

        Args:
            sample_data: Sample bytes from the file
            file_path: Path to the file for logging

        Returns:
            Best guess encoding
        """
        # Common encodings to try in order of preference
        encodings_to_try = [
            "utf-8",  # Most common modern encoding
            "utf-8-sig",  # UTF-8 with BOM
            "cp1252",  # Windows-1252 (common for Windows files)
            "iso-8859-1",  # Latin-1 (subset of Windows-1252)
            "cp1251",  # Windows-1251 (Cyrillic)
            "gb2312",  # Chinese Simplified
            "shift_jis",  # Japanese
            "euc-jp",  # Japanese
            "euc-kr",  # Korean
            "big5",  # Chinese Traditional
        ]

        for encoding in encodings_to_try:
            try:
                # Try to decode the sample
                sample_data.decode(encoding)
                logger.info(
                    f"Successfully detected encoding using fallback: {encoding}"
                )
                return encoding
            except (UnicodeDecodeError, LookupError) as e:
                logger.warning(
                    f"Failed to decode sample data with {encoding}: {str(e)}"
                )
                continue

        # If all else fails, use Latin-1 (it can decode any byte sequence)
        logger.warning(
            f"All encoding detection failed for {file_path}, using Latin-1 as last resort"
        )
        return "latin-1"

    def read_file_with_encoding_detection(self, file_path: Path) -> str:
        """
        Read a text file with automatic encoding detection.

        Args:
            file_path: Path to the text file

        Returns:
            File content as string
        """
        encoding = self.detect_file_encoding(file_path)

        try:
            with file_path.open("r", encoding=encoding, errors="replace") as f:
                content = f.read()
            logger.debug(f"Successfully read file {file_path} with encoding {encoding}")
            return content
        except Exception as e:
            logger.error(
                f"Error reading file {file_path} with encoding {encoding}: {str(e)}"
            )
            # Final fallback: read with UTF-8 and replace errors
            try:
                with file_path.open("r", encoding="utf-8", errors="replace") as f:
                    content = f.read()
                logger.warning(
                    f"Read file {file_path} with UTF-8 fallback and error replacement"
                )
                return content
            except Exception as e2:
                logger.error(
                    f"Failed to read file {file_path} even with UTF-8 fallback: {str(e2)}"
                )
                raise

    def count_file_tokens(self, file_path: str) -> int:
        """
        Count the estimated number of tokens in a file.

        Args:
            file_path: Path to the file to count tokens for

        Returns:
            Estimated number of tokens in the file
        """
        file_path = Path(file_path)
        file_extension = file_path.suffix.lower().lstrip(".")
        file_size = os.path.getsize(file_path)

        logger.info(
            f"Estimating tokens for file: {file_path}, size: {file_size/1024:.2f}KB, type: {file_extension}"
        )

        try:
            # Determine sampling parameters based on file size
            if file_size < self.SMALL_FILE_THRESHOLD:  # Less than 1MB
                sample_size = self.SAMPLE_SIZE_SMALL
                logger.debug(
                    f"Small file detected (<1MB), using sample size: {sample_size/1024:.2f}KB"
                )
                # For very small files, read the entire file
                if file_extension in ["txt", "md", "json"]:
                    logger.debug(f"Processing entire small {file_extension} file")
                    content = self.read_file_with_encoding_detection(file_path)
                    token_count = self.count_tokens(content)
                    logger.info(
                        f"Counted exact tokens for small file: {token_count} tokens"
                    )
                    return token_count
            elif file_size < self.MEDIUM_FILE_THRESHOLD:  # 1MB - 10MB
                sample_size = self.SAMPLE_SIZE_MEDIUM
                logger.debug(
                    f"Medium file detected (1-10MB), using sample size: {sample_size/1024:.2f}KB"
                )
            else:  # Larger than 10MB
                sample_size = self.SAMPLE_SIZE_LARGE
                logger.debug(
                    f"Large file detected (>10MB), using sample size: {sample_size/1024:.2f}KB"
                )

            # For larger files or specific types, use specialized estimators
            if file_extension == "pdf":
                logger.debug("Using PDF token estimator")
                token_count = self._estimate_pdf_tokens(file_path, sample_size)
                logger.info(f"Estimated PDF tokens: {token_count}")
                return token_count
            elif file_extension in ["doc", "docx"]:
                logger.debug("Using DOCX token estimator")
                token_count = self._estimate_docx_tokens(file_path, sample_size)
                logger.info(f"Estimated DOCX tokens: {token_count}")
                return token_count
            elif file_extension in ["ppt", "pptx"]:
                logger.debug("Using PPTX token estimator")
                token_count = self._estimate_pptx_tokens(file_path, sample_size)
                logger.info(f"Estimated PPTX tokens: {token_count}")
                return token_count
            elif file_extension in ["xlsx", "csv"]:
                logger.debug("Using spreadsheet token estimator")
                token_count = self._estimate_spreadsheet_tokens(file_path, sample_size)
                logger.info(f"Estimated spreadsheet tokens: {token_count}")
                return int(token_count)
            else:
                # Generic approach for other file types - use the size-based parameters
                logger.debug(
                    f"Using generic text token estimator for {file_extension} file"
                )
                token_count = self._estimate_text_tokens(
                    file_path, sample_size, self.MAX_SAMPLES
                )
                logger.info(f"Estimated generic text tokens: {token_count}")
                return int(token_count)

        except Exception as e:
            logger.error(f"Error estimating tokens for file {file_path}: {str(e)}")
            # Fall back to simple size-based estimation
            # Rough estimate: 1KB ≈ 750 tokens for text files
            token_count = int(file_size / 1024 * 750)
            logger.warning(
                f"Falling back to size-based estimation: {token_count} tokens (based on {file_size/1024:.2f}KB)"
            )
            return token_count

    def _estimate_text_tokens(
        self, file_path: Path, sample_size: int, max_samples: int
    ) -> int:
        """Estimate tokens in a text file by sampling with encoding detection."""
        file_size = os.path.getsize(file_path)

        # Detect encoding for the file
        encoding = self.detect_file_encoding(file_path)

        # For very small files, read the whole thing
        if file_size <= sample_size * 2:
            try:
                with open(
                    str(file_path), "r", encoding=encoding, errors="replace"
                ) as f:
                    content = f.read()
                return self.count_tokens(content)
            except Exception as e:
                logger.warning(
                    f"Error reading with detected encoding {encoding}, falling back to UTF-8: {str(e)}"
                )
                with open(str(file_path), "r", encoding="utf-8", errors="replace") as f:
                    content = f.read()
                return self.count_tokens(content)

        # For larger files, take samples from different positions
        samples = []
        try:
            with file_path.open("r", encoding=encoding, errors="replace") as f:
                # Start sample
                f.seek(0)
                samples.append(f.read(sample_size))

                # Middle samples
                sections = min(max_samples - 2, 3)  # -2 for start and end
                for i in range(1, sections + 1):
                    position = (file_size // (sections + 1)) * i
                    f.seek(position)
                    samples.append(f.read(sample_size))

                # End sample
                f.seek(max(0, file_size - sample_size))
                samples.append(f.read(sample_size))
        except Exception as e:
            logger.warning(
                f"Error reading with detected encoding {encoding}, falling back to UTF-8: {str(e)}"
            )
            with file_path.open("r", encoding="utf-8", errors="replace") as f:
                # Start sample
                f.seek(0)
                samples.append(f.read(sample_size))

                # Middle samples
                sections = min(max_samples - 2, 3)  # -2 for start and end
                for i in range(1, sections + 1):
                    position = (file_size // (sections + 1)) * i
                    f.seek(position)
                    samples.append(f.read(sample_size))

                # End sample
                f.seek(max(0, file_size - sample_size))
                samples.append(f.read(sample_size))

        # Count tokens in samples and extrapolate
        total_sample_size = sum(len(sample) for sample in samples)
        total_sample_tokens = sum(self.count_tokens(sample) for sample in samples)

        # Extrapolate to the full file
        if total_sample_size > 0:
            tokens_per_byte = total_sample_tokens / total_sample_size
            return int(file_size * tokens_per_byte)
        else:
            return 0

    def _estimate_pdf_tokens(self, file_path: Path, sample_size: int) -> int:
        """Estimate tokens in a PDF file using robust multi-library sampling."""

        # Validate PDF integrity before attempting token estimation
        try:
            self._validate_pdf_integrity(file_path)
        except PDFCorruptionError as e:
            logger.error(f"Cannot estimate tokens for corrupted PDF: {e.message}")
            raise e

        pdf_parsing_methods = [
            ("pypdf", self._estimate_pdf_tokens_pypdf),
            ("PyMuPDF", self._estimate_pdf_tokens_pymupdf),
            ("pdfplumber", self._estimate_pdf_tokens_pdfplumber),
        ]

        for method_name, estimate_method in pdf_parsing_methods:
            try:
                logger.debug(f"Attempting PDF token estimation with {method_name}")
                result = estimate_method(file_path, sample_size)
                logger.info(
                    f"Successfully estimated PDF tokens with {method_name}: {result}"
                )
                return result
            except Exception as e:
                logger.warning(
                    f"PDF token estimation failed with {method_name}: {str(e)}"
                )
                continue

        # Fallback to file size estimation
        logger.warning(
            "All PDF parsing methods failed for token estimation, using file size estimation"
        )
        return int(os.path.getsize(file_path) / 1024 * 500)

    def _estimate_pdf_tokens_pypdf(self, file_path: Path, sample_size: int) -> int:
        """Estimate tokens using pypdf."""
        with open(str(file_path), "rb") as file:
            reader = pypdf.PdfReader(file)
            total_pages = len(reader.pages)

            # Determine if we need to sample based on file size
            file_size = os.path.getsize(file_path)
            if file_size < self.SMALL_FILE_THRESHOLD:
                sample_pages = range(total_pages)
            else:
                sample_count = min(self.MAX_SAMPLES, total_pages)
                if sample_count <= 1:
                    sample_pages = [0]
                elif sample_count == 2:
                    sample_pages = [0, total_pages - 1]
                else:
                    sample_pages = [0]
                    step = (total_pages - 1) / (sample_count - 2)
                    for i in range(1, sample_count - 1):
                        sample_pages.append(int(i * step))
                    sample_pages.append(total_pages - 1)

            sample_text = ""

            for page_num in sample_pages:
                if page_num < total_pages:
                    page = reader.pages[page_num]
                    try:
                        sample_text += page.extract_text() + "\n\n"
                    except Exception as e:
                        error_msg = str(e)
                        logger.warning(
                            f"pypdf failed to extract text from page {page_num} for estimation: {error_msg}"
                        )
                        logger.warning(
                            "pypdf encountered an error - will try next library"
                        )
                        raise ValueError(f"pypdf cannot parse this PDF: {error_msg}")

            return self._calculate_tokens_from_sample(
                sample_text, len(sample_pages), total_pages, file_path
            )

    def _estimate_pdf_tokens_pymupdf(self, file_path: Path, sample_size: int) -> int:
        """Estimate tokens using PyMuPDF."""
        if not PYMUPDF_AVAILABLE:
            raise ImportError("PyMuPDF not available")

        doc = fitz.open(str(file_path))
        total_pages = doc.page_count

        # Determine sampling strategy
        file_size = os.path.getsize(file_path)
        if file_size < self.SMALL_FILE_THRESHOLD:
            sample_pages = range(total_pages)
        else:
            sample_count = min(self.MAX_SAMPLES, total_pages)
            if sample_count <= 1:
                sample_pages = [0]
            elif sample_count == 2:
                sample_pages = [0, total_pages - 1]
            else:
                sample_pages = [0]
                step = (total_pages - 1) / (sample_count - 2)
                for i in range(1, sample_count - 1):
                    sample_pages.append(int(i * step))
                sample_pages.append(total_pages - 1)

        sample_text = ""
        for page_num in sample_pages:
            if page_num < total_pages:
                page = doc.load_page(page_num)
                try:
                    sample_text += page.get_text() + "\n\n"
                except Exception as e:
                    logger.warning(
                        f"PyMuPDF failed to extract text from page {page_num} for estimation: {str(e)}"
                    )
                finally:
                    page = None

        doc.close()
        return self._calculate_tokens_from_sample(
            sample_text, len(sample_pages), total_pages, file_path
        )

    def _estimate_pdf_tokens_pdfplumber(self, file_path: Path, sample_size: int) -> int:
        """Estimate tokens using pdfplumber."""
        if not PDFPLUMBER_AVAILABLE:
            raise ImportError("pdfplumber not available")

        with pdfplumber.open(str(file_path)) as pdf:
            total_pages = len(pdf.pages)

            # Determine sampling strategy
            file_size = os.path.getsize(file_path)
            if file_size < self.SMALL_FILE_THRESHOLD:
                sample_pages = range(total_pages)
            else:
                sample_count = min(self.MAX_SAMPLES, total_pages)
                if sample_count <= 1:
                    sample_pages = [0]
                elif sample_count == 2:
                    sample_pages = [0, total_pages - 1]
                else:
                    sample_pages = [0]
                    step = (total_pages - 1) / (sample_count - 2)
                    for i in range(1, sample_count - 1):
                        sample_pages.append(int(i * step))
                    sample_pages.append(total_pages - 1)

            sample_text = ""
            for page_num in sample_pages:
                if page_num < total_pages:
                    page = pdf.pages[page_num]
                    try:
                        text = page.extract_text()
                        if text:
                            sample_text += text + "\n\n"
                    except Exception as e:
                        logger.warning(
                            f"pdfplumber failed to extract text from page {page_num} for estimation: {str(e)}"
                        )
                        continue

            return self._calculate_tokens_from_sample(
                sample_text, len(sample_pages), total_pages, file_path
            )

    def _calculate_tokens_from_sample(
        self,
        sample_text: str,
        sample_pages_count: int,
        total_pages: int,
        file_path: Path,
    ) -> int:
        """Calculate tokens from sampled text with fallback for scanned PDFs."""
        sample_tokens = self.count_tokens(sample_text)

        # Check if this is a scanned PDF (no extractable text)
        if sample_tokens < 10:  # Very low token count indicates scanned images
            logger.warning(
                "PDF appears to be scanned images - using realistic size-based estimation"
            )
            file_size_kb = os.path.getsize(file_path) / 1024

            # Based on typical OCR output: ~200-400 tokens per page
            estimated_tokens_per_page = min(
                400, max(100, file_size_kb / total_pages * 0.5)
            )
            estimated_tokens = int(estimated_tokens_per_page * total_pages)

            logger.info(
                f"Estimated scanned PDF tokens: {estimated_tokens} ({estimated_tokens_per_page:.0f} tokens/page)"
            )
            return estimated_tokens

        # Extrapolate to the full document
        if sample_pages_count < total_pages:
            return int(sample_tokens * (total_pages / sample_pages_count))
        else:
            return sample_tokens

    def _estimate_docx_tokens(self, file_path: Path, sample_size: int) -> int:
        """Estimate tokens in a DOCX file using size-based sampling."""
        try:
            doc = Document(file_path)
            # Sample paragraphs based on file size
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

            file_size = os.path.getsize(file_path)
            if file_size < self.SMALL_FILE_THRESHOLD:
                # For small files, process all paragraphs
                sample_paragraphs = paragraphs
            else:
                # For larger files, sample based on MAX_SAMPLES
                sample_count = min(self.MAX_SAMPLES, len(paragraphs))
                step = max(1, len(paragraphs) // sample_count)
                sample_paragraphs = [
                    paragraphs[i] for i in range(0, len(paragraphs), step)
                ][:sample_count]

            sample_text = "\n\n".join(sample_paragraphs)
            sample_tokens = self.count_tokens(sample_text)

            # Extrapolate if needed
            if len(sample_paragraphs) < len(paragraphs):
                return int(sample_tokens * (len(paragraphs) / len(sample_paragraphs)))
            else:
                return sample_tokens

        except Exception as e:
            logger.error(f"Error estimating DOCX tokens: {str(e)}")
            # Fall back to size-based estimation
            return int(
                os.path.getsize(file_path) / 1024 * 600
            )  # Rough estimate for DOCX

    def _estimate_spreadsheet_tokens(self, file_path: Path, sample_size: int) -> int:
        """Estimate tokens in a spreadsheet file using size-based sampling with encoding detection."""
        try:
            # Determine number of rows to sample based on file size
            file_size = os.path.getsize(file_path)
            if file_size < self.SMALL_FILE_THRESHOLD:
                sample_rows = 100  # Default for small files
            elif file_size < self.MEDIUM_FILE_THRESHOLD:
                sample_rows = 50  # Fewer for medium files to avoid memory issues
            else:
                sample_rows = 20  # Even fewer for large files

            if file_path.suffix.lower() == ".csv":
                # Detect encoding for CSV files
                encoding = self.detect_file_encoding(file_path)
                logger.debug(f"Reading CSV with detected encoding: {encoding}")

                try:
                    df = pd.read_csv(
                        file_path,
                        nrows=sample_rows,
                        encoding=encoding,
                        low_memory=False,
                    )
                except Exception as e:
                    logger.warning(
                        f"Failed to read CSV with detected encoding {encoding}, trying fallbacks: {str(e)}"
                    )
                    # Try common encodings as fallback
                    for fallback_encoding in ["utf-8", "latin-1", "cp1252"]:
                        try:
                            df = pd.read_csv(
                                file_path,
                                nrows=sample_rows,
                                encoding=fallback_encoding,
                                low_memory=False,
                            )
                            logger.info(
                                f"Successfully read CSV with fallback encoding: {fallback_encoding}"
                            )
                            encoding = fallback_encoding
                            break
                        except Exception as e:
                            logger.warning(
                                f"Failed to read CSV with fallback encoding {fallback_encoding}: {str(e)}"
                            )
                            continue
                    else:
                        # Last resort: let pandas auto-detect
                        df = pd.read_csv(
                            file_path,
                            nrows=sample_rows,
                            encoding=None,
                            low_memory=False,
                        )
                        encoding = "auto-detected"

                file_type = "csv"
            else:  # Excel
                # For Excel, read the first sheet with sample rows
                df = pd.read_excel(file_path, sheet_name=0, nrows=sample_rows)
                file_type = "xlsx"
                encoding = "excel-binary"

            # Convert sample to string and count tokens
            sample_text = df.to_string(index=False)
            sample_tokens = self.count_tokens(sample_text)

            # Estimate total rows and extrapolate
            if file_type == "csv":
                # Estimate total rows in CSV using the detected encoding
                try:
                    with open(
                        str(file_path), "r", encoding=encoding, errors="replace"
                    ) as f:
                        total_rows = sum(1 for _ in f)
                except Exception as e:
                    logger.warning(
                        f"Error counting rows with encoding {encoding}: {str(e)}"
                    )
                    # Fallback to UTF-8 with error replacement
                    with open(
                        str(file_path), "r", encoding="utf-8", errors="replace"
                    ) as f:
                        total_rows = sum(1 for _ in f)
            else:
                # For Excel, this is more complex - use file size as a proxy
                sample_size = len(sample_text.encode("utf-8"))
                # Rough estimate based on file size ratio
                total_rows = max(sample_rows, int((file_size / sample_size) * len(df)))

            # Extrapolate tokens based on row count ratio
            if len(df) > 0:
                return int(sample_tokens * (total_rows / len(df)))
            else:
                return 0

        except Exception as e:
            logger.error(f"Error estimating spreadsheet tokens: {str(e)}")
            # Fall back to size-based estimation
            return int(
                os.path.getsize(file_path) / 1024 * 400
            )  # Rough estimate for spreadsheets

    def _estimate_pptx_tokens(self, file_path: Path, sample_size: int) -> int:
        """Estimate tokens in a PPTX file using size-based sampling."""
        try:
            prs = Presentation(file_path)
            # Extract text from all slides
            slide_texts = []

            for slide in prs.slides:
                slide_text = ""
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"

                if slide_text.strip():
                    slide_texts.append(slide_text.strip())

            file_size = os.path.getsize(file_path)
            if file_size < self.SMALL_FILE_THRESHOLD:
                # For small files, process all slides
                sample_slides = slide_texts
            else:
                # For larger files, sample based on MAX_SAMPLES
                sample_count = min(self.MAX_SAMPLES, len(slide_texts))
                if sample_count > 0:
                    step = max(1, len(slide_texts) // sample_count)
                    sample_slides = [
                        slide_texts[i] for i in range(0, len(slide_texts), step)
                    ][:sample_count]
                else:
                    sample_slides = slide_texts

            sample_text = "\n\n".join(sample_slides)
            sample_tokens = self.count_tokens(sample_text)

            # Extrapolate if needed
            if len(sample_slides) < len(slide_texts) and len(sample_slides) > 0:
                return int(sample_tokens * (len(slide_texts) / len(sample_slides)))
            else:
                return sample_tokens

        except Exception as e:
            logger.error(f"Error estimating PPTX tokens: {str(e)}")
            # Fall back to size-based estimation
            return int(
                os.path.getsize(file_path) / 1024 * 500
            )  # Rough estimate for PPTX

    def split_file(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Split a file into smaller chunks based on its type.

        Args:
            file_path: Path to the file to split

        Returns:
            List of dictionaries containing information about each split
        """
        file_path = Path(file_path)
        file_extension = file_path.suffix.lower().lstrip(".")

        logger.info(
            f"Splitting file: {file_path}, type: {file_extension}"
        )  # Choose the appropriate splitting method based on file type
        if file_extension in ["txt", "md", "json"]:
            return self._split_text_file(file_path)
        elif file_extension == "pdf":
            return self._split_pdf_file(file_path)
        elif file_extension in ["doc", "docx"]:
            return self._split_docx_file(file_path)
        elif file_extension in ["ppt", "pptx"]:
            return self._split_pptx_file(file_path)
        elif file_extension in ["xlsx", "csv"]:
            return self._split_spreadsheet_file(file_path)
        elif file_extension in ["html", "xml"]:
            return self._split_markup_file(file_path)
        else:
            logger.warning(
                f"Unsupported file type for splitting: {file_extension}. Falling back to generic text splitting."
            )
            return self._split_text_file(file_path)

    def _split_text_file(self, file_path: Path) -> List[Dict[str, Any]]:
        """Split text-based files (txt, md, json, etc.) with encoding detection."""
        try:
            # Use encoding detection for text files
            content = self.read_file_with_encoding_detection(file_path)

            # Special handling for JSON files to preserve structure
            if file_path.suffix.lower() == ".json":
                logger.debug(f"Processing JSON file: {file_path}")
                try:
                    # Parse the JSON to validate and preserve structure
                    json_content = json.loads(content)

                    # If it's an array, we can split by items
                    if isinstance(json_content, list):
                        logger.debug(f"JSON is a list with {len(json_content)} items")
                        # Calculate approximately how many items should go in each split
                        total_tokens = self.count_tokens(content)
                        if total_tokens <= self.max_tokens_per_split:
                            # If the entire file fits within token limit, keep it as one chunk
                            return self._create_splits_from_chunks(
                                [content], str(file_path), file_type="json"
                            )

                        # Estimate tokens per item
                        avg_tokens_per_item = total_tokens / len(json_content)
                        items_per_split = max(
                            1, int(self.max_tokens_per_split / avg_tokens_per_item)
                        )

                        # Split the array into chunks
                        chunks = []
                        for i in range(0, len(json_content), items_per_split):
                            chunk = json_content[i : i + items_per_split]
                            # Format as proper JSON
                            chunk_content = json.dumps(chunk)
                            chunks.append(chunk_content)

                        logger.debug(f"Split JSON array into {len(chunks)} chunks")
                        return self._create_splits_from_chunks(
                            chunks, str(file_path), file_type="json"
                        )
                    else:
                        # For non-array JSON, keep intact
                        logger.debug("JSON is not an array, keeping intact")
                        return self._create_splits_from_chunks(
                            [content], str(file_path), file_type="json"
                        )
                except json.JSONDecodeError as e:
                    logger.warning(
                        f"Invalid JSON in {file_path}: {e}. Falling back to text splitting."
                    )

            # For Markdown, consider rendering to HTML first to preserve structure
            if file_path.suffix.lower() == ".md":
                content_html = markdown.markdown(content)
                soup = BeautifulSoup(content_html, "html.parser")
                content = soup.get_text()

            # Split the content into logical chunks (paragraphs)
            paragraphs = [p for p in content.split("\n\n") if p.strip()]

            logger.debug(f"Split text file into {len(paragraphs)} paragraphs")
            return self._create_splits_from_chunks(
                paragraphs, str(file_path), file_type="text"
            )

        except Exception as e:
            logger.error(f"Error splitting text file {file_path}: {str(e)}")
            raise

    def _open_pdf_reader(self, file_path: Path) -> Tuple[pypdf.PdfReader, int]:
        """Open PDF file and return reader with total pages count."""
        with open(str(file_path), "rb") as file:
            reader = pypdf.PdfReader(file)
            total_pages = len(reader.pages)
            return reader, total_pages

    def _extract_pdf_content(
        self, reader: pypdf.PdfReader, total_pages: int
    ) -> Tuple[List[Tuple[str, int]], int]:
        """Extract textual content from PDF pages."""
        pages = []
        total_text_length = 0

        for page_num in range(total_pages):
            page = reader.pages[page_num]
            text = page.extract_text()
            if text and text.strip():
                pages.append((text, page_num))
                total_text_length += len(text)
            else:
                # For pages with no extractable text (scanned), add placeholder
                pages.append((f"[Scanned page {page_num + 1}]", page_num))

        return pages, total_text_length

    def _log_pdf_content_info(
        self, pages: List[Tuple[str, int]], total_pages: int
    ) -> None:
        """Log information about PDF content and extractable text."""
        extractable_pages = len([p for p, _ in pages if not p.startswith("[Scanned")])
        logger.info(
            f"PDF has {total_pages} total pages, {extractable_pages} with extractable text"
        )

    def _is_scanned_pdf(
        self, pages: List[Tuple[str, int]], total_text_length: int
    ) -> bool:
        """Determine if PDF is entirely scanned (no extractable text)."""
        extractable_pages = len([p for p, _ in pages if not p.startswith("[Scanned")])
        return total_text_length == 0 or extractable_pages == 0

    def _process_text_based_pdf(
        self, pages: List[Tuple[str, int]], file_path: Path
    ) -> List[Dict[str, Any]]:
        """Process PDF with extractable text content."""
        logger.debug(f"Split PDF file into {len(pages)} pages with text content")
        return self._create_splits_from_chunks(
            [page_text for page_text, _ in pages],
            str(file_path),
            file_type="pdf",
            page_numbers=[page_num for _, page_num in pages],
            original_pdf=file_path,
        )

    def _split_pdf_file(self, file_path: Path) -> List[Dict[str, Any]]:
        """Split PDF files with robust parsing using multiple libraries as fallbacks."""

        # Step 1: Validate PDF integrity before attempting any processing
        try:
            self._validate_pdf_integrity(file_path)
        except PDFCorruptionError as e:
            logger.error(f"PDF corruption detected: {e.message}")
            if e.corruption_details:
                logger.error(f"Corruption details: {e.corruption_details}")
            raise e

        pdf_parsing_methods = [
            ("pypdf", self._extract_text_with_pypdf),
            ("PyMuPDF", self._extract_text_with_pymupdf),
            ("pdfplumber", self._extract_text_with_pdfplumber),
        ]

        last_error = None

        for method_name, extract_method in pdf_parsing_methods:
            try:
                logger.debug(f"Attempting PDF parsing with {method_name}")
                pages, total_text_length = extract_method(file_path)

                if pages:
                    logger.info(
                        f"Successfully parsed PDF with {method_name} - {len(pages)} pages, {total_text_length} characters"
                    )
                    # Step 3: Log information about the PDF content
                    self._log_pdf_content_info(pages, len(pages))

                    # Step 4: Determine processing strategy based on content
                    if self._is_scanned_pdf(pages, total_text_length):
                        logger.warning(
                            f"PDF appears to be entirely scanned images - creating page-based splits using {method_name}"
                        )
                        return self._create_scanned_pdf_splits(file_path, len(pages))

                    # Step 5: Process PDF with extractable text content
                    return self._process_text_based_pdf(pages, file_path)

            except Exception as e:
                last_error = e
                logger.warning(f"PDF parsing failed with {method_name}: {str(e)}")
                continue

        # If all methods failed, raise the last error
        if last_error:
            logger.error(
                f"All PDF parsing methods failed. Last error: {str(last_error)}"
            )
            raise last_error
        else:
            raise Exception("Unable to parse PDF with any available method")

    def _extract_text_with_pypdf(
        self, file_path: Path
    ) -> Tuple[List[Tuple[str, int]], int]:
        """Extract text using pypdf (original method)."""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PDF file not found: {file_path}")

        pages = []
        total_text_length = 0

        with open(str(file_path), "rb") as file:
            reader = pypdf.PdfReader(file)
            total_pages = len(reader.pages)

            for page_num in range(total_pages):
                page = reader.pages[page_num]
                try:
                    text = page.extract_text()
                    if text and text.strip():
                        pages.append((text, page_num))
                        total_text_length += len(text)
                    else:
                        # For pages with no extractable text (scanned), add placeholder
                        pages.append((f"[Scanned page {page_num + 1}]", page_num))
                except Exception as e:
                    error_msg = str(e)
                    logger.warning(
                        f"pypdf failed to extract text from page {page_num}: {error_msg}"
                    )
                    logger.warning("pypdf encountered an error - will try next library")
                    raise ValueError(f"pypdf cannot parse this PDF: {error_msg}")

        return pages, total_text_length

    def _extract_text_with_pymupdf(
        self, file_path: Path
    ) -> Tuple[List[Tuple[str, int]], int]:
        """Extract text using PyMuPDF (fitz) - more robust for problematic PDFs."""
        if not PYMUPDF_AVAILABLE:
            raise ImportError("PyMuPDF (fitz) not available")

        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PDF file not found: {file_path}")

        pages = []
        total_text_length = 0

        try:
            doc = fitz.open(str(file_path))
            total_pages = doc.page_count

            for page_num in range(total_pages):
                page = doc.load_page(page_num)
                try:
                    text = page.get_text()
                    if text and text.strip():
                        pages.append((text, page_num))
                        total_text_length += len(text)
                    else:
                        # For pages with no extractable text (scanned), add placeholder
                        pages.append((f"[Scanned page {page_num + 1}]", page_num))
                except Exception as e:
                    logger.warning(
                        f"PyMuPDF failed to extract text from page {page_num}: {str(e)}"
                    )
                    pages.append((f"[Error extracting page {page_num + 1}]", page_num))
                finally:
                    page = None  # Free page resources

            doc.close()
            return pages, total_text_length

        except Exception as e:
            logger.error(f"PyMuPDF failed to open PDF: {str(e)}")
            raise

    def _extract_text_with_pdfplumber(
        self, file_path: Path
    ) -> Tuple[List[Tuple[str, int]], int]:
        """Extract text using pdfplumber - good for structured PDFs."""
        if not PDFPLUMBER_AVAILABLE:
            raise ImportError("pdfplumber not available")

        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PDF file not found: {file_path}")

        pages = []
        total_text_length = 0

        try:
            with pdfplumber.open(str(file_path)) as pdf:
                total_pages = len(pdf.pages)

                for page_num in range(total_pages):
                    page = pdf.pages[page_num]
                    try:
                        text = page.extract_text()
                        if text and text.strip():
                            pages.append((text, page_num))
                            total_text_length += len(text)
                        else:
                            # For pages with no extractable text (scanned), add placeholder
                            pages.append((f"[Scanned page {page_num + 1}]", page_num))
                    except Exception as e:
                        logger.warning(
                            f"pdfplumber failed to extract text from page {page_num}: {str(e)}"
                        )
                        pages.append(
                            (f"[Error extracting page {page_num + 1}]", page_num)
                        )

            return pages, total_text_length

        except Exception as e:
            logger.error(f"pdfplumber failed to open PDF: {str(e)}")
            raise

    def _create_scanned_pdf_splits(
        self, file_path: Path, total_pages: int
    ) -> List[Dict[str, Any]]:
        """Create splits for scanned PDFs based on page ranges."""
        try:
            # Calculate pages per split based on estimated tokens per page
            file_size_kb = os.path.getsize(file_path) / 1024

            # Based on typical OCR output: ~200-400 tokens per page
            estimated_tokens_per_page = min(
                400, max(100, file_size_kb / total_pages * 0.5)
            )
            pages_per_split = max(
                1, int(self.max_tokens_per_split / estimated_tokens_per_page)
            )

            logger.info(
                f"Creating scanned PDF splits: ~{estimated_tokens_per_page:.0f} tokens/page, {pages_per_split} pages/split"
            )

            splits = []

            for start_page in range(0, total_pages, pages_per_split):
                end_page = min(start_page + pages_per_split, total_pages)
                page_range = list(range(start_page, end_page))

                # Estimate tokens for this split
                split_token_count = int(estimated_tokens_per_page * len(page_range))

                # Create the actual PDF split file
                split_path = self._generate_split_path(str(file_path))
                self._create_pdf_page_range_split(split_path, file_path, page_range)

                # Get the file size
                split_size = os.path.getsize(split_path)

                splits.append(
                    {
                        "split_file_path": split_path,
                        "token_count": split_token_count,
                        "size": split_size,
                        "split_index": len(splits),
                        "total_splits": -1,  # Will be updated later
                        "page_range": f"{start_page + 1}-{end_page}",
                        "is_scanned": True,
                    }
                )

                logger.info(
                    f"Created scanned PDF split {len(splits)}: pages {start_page + 1}-{end_page} ({split_token_count} tokens)"
                )

            # Update total splits count
            for split in splits:
                split["total_splits"] = len(splits)

            logger.info(f"Created {len(splits)} splits for scanned PDF")
            return splits

        except Exception as e:
            logger.error(f"Error creating scanned PDF splits: {str(e)}")
            raise

    def _create_pdf_page_range_split(
        self, split_path: str, original_pdf: Path, page_range: List[int]
    ) -> None:
        """Create a PDF split containing specific page ranges using PyMuPDF (robust fallback to pypdf)."""
        # Try PyMuPDF first (more robust)
        if PYMUPDF_AVAILABLE:
            try:
                self._create_pdf_page_range_split_pymupdf(
                    split_path, original_pdf, page_range
                )
                logger.debug(
                    f"Successfully created PDF split with PyMuPDF: pages {page_range}"
                )
                return
            except Exception as e:
                logger.warning(
                    f"PyMuPDF PDF split creation failed: {str(e)}, trying pypdf fallback"
                )

        # Fallback to pypdf
        try:
            with open(original_pdf, "rb") as file:
                reader = pypdf.PdfReader(file)
                writer = pypdf.PdfWriter()

                # Add each page in the range to the new PDF
                for page_num in page_range:
                    if page_num < len(reader.pages):
                        writer.add_page(reader.pages[page_num])
                    else:
                        logger.warning(
                            f"Page {page_num} is out of range for PDF with {len(reader.pages)} pages"
                        )

                # Save the new PDF
                with open(split_path, "wb") as output_file:
                    writer.write(output_file)

                logger.debug(
                    f"Successfully created PDF split with pypdf: pages {page_range} at: {split_path}"
                )

        except Exception as e:
            logger.error(f"Both PyMuPDF and pypdf PDF split creation failed: {e}")
            raise

    def _create_pdf_page_range_split_pymupdf(
        self, split_path: str, original_pdf: Path, page_range: List[int]
    ) -> None:
        """Create a PDF split using PyMuPDF (more robust for corrupted PDFs)."""
        try:
            # Open source document
            src_doc = fitz.open(str(original_pdf))

            # Create new document for the split
            dst_doc = fitz.open()

            # Copy specified pages to new document
            for page_num in page_range:
                if page_num < src_doc.page_count:
                    dst_doc.insert_pdf(src_doc, from_page=page_num, to_page=page_num)
                else:
                    logger.warning(
                        f"Page {page_num} is out of range for PDF with {src_doc.page_count} pages"
                    )

            # Save the new PDF
            dst_doc.save(split_path)
            dst_doc.close()
            src_doc.close()

            logger.debug(
                f"PyMuPDF created PDF split: {len(page_range)} pages at {split_path}"
            )

        except Exception as e:
            logger.error(f"PyMuPDF PDF split creation error: {e}")
            raise

    def _split_docx_file(self, file_path: Path) -> List[Dict[str, Any]]:
        """Split DOCX files with consideration for section boundaries."""
        try:
            doc = Document(file_path)

            # Extract paragraphs from the document
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

            logger.debug(f"Split DOCX file into {len(paragraphs)} paragraphs")

            return self._create_splits_from_chunks(
                paragraphs, str(file_path), file_type="docx"
            )

        except Exception as e:
            logger.error(f"Error splitting DOCX file {file_path}: {str(e)}")
            raise

    def _split_pptx_file(self, file_path: Path) -> List[Dict[str, Any]]:
        """Split PPTX files with consideration for slide boundaries."""
        try:
            prs = Presentation(file_path)

            # Extract text from each slide
            slide_contents = []

            for slide_num, slide in enumerate(prs.slides):
                slide_text = ""

                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"

                # Only add slides with meaningful content
                if slide_text.strip():
                    slide_contents.append(
                        f"Slide {slide_num + 1}:\n{slide_text.strip()}"
                    )
                else:
                    # Add placeholder for slides without text (might have images/charts)
                    slide_contents.append(
                        f"Slide {slide_num + 1}: [Slide contains visual content]"
                    )

            logger.debug(f"Split PPTX file into {len(slide_contents)} slides")

            return self._create_splits_from_chunks(
                slide_contents, str(file_path), file_type="pptx"
            )

        except Exception as e:
            logger.error(f"Error splitting PPTX file {file_path}: {str(e)}")
            raise

    def _split_spreadsheet_file(self, file_path: Path) -> List[Dict[str, Any]]:
        """Split Excel or CSV files by worksheet or chunks of rows with encoding detection."""
        try:
            if file_path.suffix.lower() == ".csv":
                # Detect encoding for CSV files
                encoding = self.detect_file_encoding(file_path)
                logger.debug(
                    f"Reading CSV for splitting with detected encoding: {encoding}"
                )

                try:
                    df = pd.read_csv(file_path, encoding=encoding, low_memory=False)
                except Exception as e:
                    logger.warning(
                        f"Failed to read CSV with detected encoding {encoding}, trying fallbacks: {str(e)}"
                    )
                    # Try common encodings as fallback
                    for fallback_encoding in ["utf-8", "latin-1", "cp1252"]:
                        try:
                            df = pd.read_csv(
                                file_path, encoding=fallback_encoding, low_memory=False
                            )
                            logger.info(
                                f"Successfully read CSV with fallback encoding: {fallback_encoding}"
                            )
                            encoding = fallback_encoding
                            break
                        except Exception as e:
                            logger.warning(
                                f"Failed to read CSV with fallback encoding {fallback_encoding}: {str(e)}"
                            )
                            continue
                    else:
                        # Last resort: let pandas auto-detect
                        df = pd.read_csv(file_path, encoding=None, low_memory=False)
                        encoding = "auto-detected"

                sheets = {"Sheet1": df}
                file_type = "csv"
            else:  # Excel
                sheets = pd.read_excel(file_path, sheet_name=None)
                file_type = "xlsx"

            chunks = []
            sheet_info = []

            # Process each sheet
            for sheet_name, df in sheets.items():
                # Convert dataframe to string representation
                sheet_text = f"Sheet: {sheet_name}\n\n"
                sheet_text += df.to_string(index=False)

                # Split large sheets into chunks based on rows
                if self.count_tokens(sheet_text) > self.max_tokens_per_split:
                    rows_per_chunk = max(
                        5,
                        len(df)
                        // (
                            self.count_tokens(sheet_text) // self.max_tokens_per_split
                            + 1
                        ),
                    )
                    for i in range(0, len(df), rows_per_chunk):
                        chunk_df = df.iloc[i : i + rows_per_chunk]
                        chunk_text = f"Sheet: {sheet_name} (Rows {i+1}-{min(i+rows_per_chunk, len(df))})\n\n"
                        chunk_text += chunk_df.to_string(index=False)
                        chunks.append(chunk_text)
                        # Store sheet name and row range for proper Excel splitting
                        sheet_info.append(
                            (sheet_name, i, min(i + rows_per_chunk, len(df)))
                        )
                else:
                    chunks.append(sheet_text)
                    sheet_info.append((sheet_name, 0, len(df)))

            logger.debug(f"Split spreadsheet file into {len(chunks)} chunks")

            return self._create_splits_from_chunks(
                chunks,
                str(file_path),
                file_type=file_type,
                sheet_info=sheet_info,
                original_spreadsheet=file_path,
            )

        except Exception as e:
            logger.error(f"Error splitting spreadsheet file {file_path}: {str(e)}")
            raise

    def _split_markup_file(self, file_path: Path) -> List[Dict[str, Any]]:
        """Split HTML or XML files by elements with encoding detection."""
        try:
            # Use encoding detection for markup files
            content = self.read_file_with_encoding_detection(file_path)

            file_type = file_path.suffix.lower().lstrip(".")

            if file_type == "html":
                sections, tags = self._extract_html_sections(content)
                html_content = tags
            else:  # XML
                sections, xml_content = self._extract_xml_sections(content)
                html_content = xml_content

            logger.debug(f"Split markup file into {len(sections)} sections")

            return self._create_splits_from_chunks(
                sections,
                str(file_path),
                file_type=file_type,
                markup_content=html_content,
            )

        except Exception as e:
            logger.error(f"Error splitting markup file {file_path}: {str(e)}")
            raise

    def _extract_html_sections(self, content: str) -> Tuple[List[str], List[str]]:
        """Extract sections and tags from HTML content."""
        soup = BeautifulSoup(content, "html.parser")

        # Extract text by sections (div, section, article, etc.)
        sections = []
        tags = []

        # Try to find structural elements first
        for tag in ["article", "section", "div", "main"]:
            elements = soup.find_all(tag)
            for element in elements:
                if element.get_text().strip():
                    sections.append(element.get_text())
                    tags.append(f"<{tag}>{element.get_text()}</{tag}>")

        # If no sections found, fall back to paragraphs
        if not sections:
            sections = [
                p.get_text() for p in soup.find_all("p") if p.get_text().strip()
            ]
            tags = [
                f"<p>{p.get_text()}</p>"
                for p in soup.find_all("p")
                if p.get_text().strip()
            ]

        # If still no content, use the whole document
        if not sections:
            sections = [soup.get_text()]
            tags = [str(soup)]

        return sections, tags

    def _extract_xml_sections(self, content: str) -> Tuple[List[str], List[str]]:
        """Extract sections and tags from XML content."""
        root = ET.fromstring(content)  # using the secure defusedxml version

        # Process XML by extracting text from each element
        sections = []
        xml_content = []

        for element in root.iter():
            if element.text and element.text.strip():
                sections.append(f"{element.tag}: {element.text.strip()}")
                element_str = ET.tostring(element, encoding="unicode")
                xml_content.append(element_str)

        return sections, xml_content

    def _create_splits_from_chunks(
        self,
        chunks: List[str],
        original_file_path: str,
        file_type: str = "text",
        page_numbers: List[int] = None,
        sheet_info: List[Tuple[str, int, int]] = None,
        markup_content: List[str] = None,
        original_pdf: Path = None,
        original_spreadsheet: Path = None,
    ) -> List[Dict[str, Any]]:
        """
        Create file splits from chunks of content.

        Args:
            chunks: List of content chunks (paragraphs, pages, etc.)
            original_file_path: Path to the original file
            file_type: Type of the file being split
            page_numbers: List of page numbers for PDF files
            sheet_info: List of sheet information for spreadsheet files
            markup_content: List of markup elements for HTML/XML files
            original_pdf: Path to the original PDF file
            original_spreadsheet: Path to the original spreadsheet file

        Returns:
            List of dictionaries with information about each split
        """
        splits = []
        current_split_content = ""
        current_token_count = 0
        current_chunk_indices = []

        # Process each chunk
        for i, chunk in enumerate(chunks):
            chunk_token_count = self.count_tokens(chunk)

            # If a single chunk exceeds max tokens, handle it separately
            if chunk_token_count > self.max_tokens_per_split:
                # Save any accumulated content before processing the large chunk
                if current_token_count > 0:
                    split_info = self._create_split_info(
                        current_split_content,
                        current_token_count,
                        original_file_path,
                        file_type,
                        current_chunk_indices,
                        page_numbers,
                        sheet_info,
                        markup_content,
                        original_pdf,
                        original_spreadsheet,
                    )
                    splits.append(split_info)
                    current_split_content = ""
                    current_token_count = 0
                    current_chunk_indices = []

                # Process large chunk by breaking it into sentences
                sentence_splits = self._process_large_chunk(
                    chunk,
                    i,
                    original_file_path,
                    file_type,
                    page_numbers,
                    sheet_info,
                    markup_content,
                    original_pdf,
                    original_spreadsheet,
                )
                splits.extend(sentence_splits)
            else:
                # Handle normal-sized chunks
                result = self._add_chunk_to_current_split(
                    chunk,
                    chunk_token_count,
                    i,
                    current_split_content,
                    current_token_count,
                    current_chunk_indices,
                    original_file_path,
                    file_type,
                    page_numbers,
                    sheet_info,
                    markup_content,
                    original_pdf,
                    original_spreadsheet,
                )

                # Unpack the result
                (
                    splits_update,
                    current_split_content,
                    current_token_count,
                    current_chunk_indices,
                ) = result
                if splits_update:
                    splits.append(splits_update)

        # Save any remaining content
        if current_token_count > 0:
            split_info = self._create_split_info(
                current_split_content,
                current_token_count,
                original_file_path,
                file_type,
                current_chunk_indices,
                page_numbers,
                sheet_info,
                markup_content,
                original_pdf,
                original_spreadsheet,
            )
            splits.append(split_info)

        # Update the total_splits count in each split
        return self._update_split_indices(splits)

    def _process_large_chunk(
        self,
        chunk: str,
        chunk_index: int,
        original_file_path: str,
        file_type: str,
        page_numbers: List[int] = None,
        sheet_info: List[Tuple[str, int, int]] = None,
        markup_content: List[str] = None,
        original_pdf: Path = None,
        original_spreadsheet: Path = None,
    ) -> List[Dict[str, Any]]:
        """Process a chunk that exceeds the maximum token limit by splitting it into sentences."""
        splits = []
        sentences = self._split_into_sentences(chunk)
        temp_content = ""
        temp_token_count = 0
        temp_chunk_indices = []

        for sentence in sentences:
            sentence_token_count = self.count_tokens(sentence)

            # Check if adding this sentence would exceed the limit and we have enough content
            if (
                temp_token_count + sentence_token_count > self.max_tokens_per_split
                and temp_token_count >= self.min_tokens_per_split
            ):
                # Save current accumulated content
                split_info = self._create_split_info(
                    temp_content,
                    temp_token_count,
                    original_file_path,
                    file_type,
                    [chunk_index],
                    page_numbers,
                    sheet_info,
                    markup_content,
                    original_pdf,
                    original_spreadsheet,
                )
                splits.append(split_info)

                # Start a new split with this sentence
                temp_content = sentence
                temp_token_count = sentence_token_count
                temp_chunk_indices = [chunk_index]
            else:
                # Add sentence to current content
                if temp_content:
                    temp_content += " " + sentence
                else:
                    temp_content = sentence
                temp_token_count += sentence_token_count
                if chunk_index not in temp_chunk_indices:
                    temp_chunk_indices.append(chunk_index)

        # Save any remaining content from sentences
        if temp_token_count > 0:
            split_info = self._create_split_info(
                temp_content,
                temp_token_count,
                original_file_path,
                file_type,
                temp_chunk_indices,
                page_numbers,
                sheet_info,
                markup_content,
                original_pdf,
                original_spreadsheet,
            )
            splits.append(split_info)

        return splits

    def _add_chunk_to_current_split(
        self,
        chunk: str,
        chunk_token_count: int,
        chunk_index: int,
        current_split_content: str,
        current_token_count: int,
        current_chunk_indices: List[int],
        original_file_path: str,
        file_type: str,
        page_numbers: List[int] = None,
        sheet_info: List[Tuple[str, int, int]] = None,
        markup_content: List[str] = None,
        original_pdf: Path = None,
        original_spreadsheet: Path = None,
    ) -> Tuple[Optional[Dict[str, Any]], str, int, List[int]]:
        """Add a chunk to the current split or start a new one if needed."""
        new_split = None

        # Check if adding this chunk would exceed the limit and we have enough content
        if (
            current_token_count + chunk_token_count > self.max_tokens_per_split
            and current_token_count >= self.min_tokens_per_split
        ):
            # Save current accumulated content
            new_split = self._create_split_info(
                current_split_content,
                current_token_count,
                original_file_path,
                file_type,
                current_chunk_indices,
                page_numbers,
                sheet_info,
                markup_content,
                original_pdf,
                original_spreadsheet,
            )

            # Start a new split with this chunk
            current_split_content = chunk
            current_token_count = chunk_token_count
            current_chunk_indices = [chunk_index]
        else:
            # Add chunk to current content
            if current_split_content:
                current_split_content += "\n\n" + chunk
            else:
                current_split_content = chunk
            current_token_count += chunk_token_count
            current_chunk_indices.append(chunk_index)

        return (
            new_split,
            current_split_content,
            current_token_count,
            current_chunk_indices,
        )

    def _create_split_info(
        self,
        content: str,
        token_count: int,
        original_file_path: str,
        file_type: str,
        chunk_indices: List[int],
        page_numbers: List[int] = None,
        sheet_info: List[Tuple[str, int, int]] = None,
        markup_content: List[str] = None,
        original_pdf: Path = None,
        original_spreadsheet: Path = None,
    ) -> Dict[str, Any]:
        """Create split information by saving content to file and returning metadata."""
        return self._save_split(
            content,
            token_count,
            original_file_path,
            file_type,
            chunk_indices,
            page_numbers,
            sheet_info,
            markup_content,
            original_pdf,
            original_spreadsheet,
        )

    def _update_split_indices(
        self, splits: List[Dict[str, Any]]
    ) -> List[Dict[str, Any]]:
        """Update split indices and total count."""
        total_splits = len(splits)
        logger.info(f"Finalizing {total_splits} file splits")
        for i, split in enumerate(splits):
            split["split_index"] = i
            split["total_splits"] = total_splits
        return splits

    def _split_into_sentences(self, text: str) -> List[str]:
        """Split text into sentences for finer-grained chunking."""
        # Simple sentence splitting - can be improved with NLP libraries
        sentences = []
        for para in text.split("\n"):
            for sentence in para.replace("!", ".").replace("?", ".").split("."):
                if sentence.strip():
                    stripped_sentence = sentence.strip()
                    if not stripped_sentence.endswith("."):
                        stripped_sentence += "."
                    sentences.append(stripped_sentence)
        logger.debug(f"Split text into {len(sentences)} sentences")
        return sentences

    def _save_split(
        self,
        content: str,
        token_count: int,
        original_file_path: str,
        file_type: str = "text",
        chunk_indices: List[int] = None,
        page_numbers: List[int] = None,
        sheet_info: List[Tuple[str, int, int]] = None,
        markup_content: List[str] = None,
        original_pdf: Path = None,
        original_spreadsheet: Path = None,
    ) -> Dict[str, Any]:
        """
        Save split content to a file and return information about the split.

        Args:
            content: The split content
            token_count: Number of tokens in the content
            original_file_path: Path to the original file
            file_type: Type of the file being split
            chunk_indices: Indices of the chunks included in this split
            page_numbers: List of page numbers for PDF files
            sheet_info: List of sheet information for spreadsheet files
            markup_content: List of markup elements for HTML/XML files
            original_pdf: Path to the original PDF file
            original_spreadsheet: Path to the original spreadsheet file

        Returns:
            Dictionary with information about the split
        """
        # Generate a unique filename for this split
        split_path = self._generate_split_path(original_file_path)
        logger.debug(
            f"Saving split to {split_path} with {token_count} tokens, file type: {file_type}"
        )

        # Dispatch to the appropriate handler based on file type
        self._save_split_by_type(
            file_type=file_type,
            split_path=split_path,
            content=content,
            original_pdf=original_pdf,
            page_numbers=page_numbers,
            chunk_indices=chunk_indices,
            original_spreadsheet=original_spreadsheet,
            sheet_info=sheet_info,
            markup_content=markup_content,
            original_file_path=original_file_path,
        )

        # Get the file size
        split_size = os.path.getsize(split_path)
        logger.debug(f"Split file size: {split_size/1024:.2f}KB")

        return {
            "split_file_path": split_path,
            "token_count": token_count,
            "size": split_size,
            "split_index": -1,  # Will be updated later
            "total_splits": -1,  # Will be updated later
        }

    def _save_split_by_type(
        self,
        file_type: str,
        split_path: str,
        content: str,
        original_pdf: Path = None,
        page_numbers: List[int] = None,
        chunk_indices: List[int] = None,
        original_spreadsheet: Path = None,
        sheet_info: List[Tuple[str, int, int]] = None,
        markup_content: List[str] = None,
        original_file_path: str = None,
    ) -> None:
        """
        Save split content based on file type.

        Args:
            file_type: Type of the file being split
            split_path: Path where to save the split
            content: The content to save
            original_pdf: Path to the original PDF file
            page_numbers: List of page numbers for PDF files
            chunk_indices: Indices of the chunks included in this split
            original_spreadsheet: Path to the original spreadsheet file
            sheet_info: List of sheet information for spreadsheet files
            markup_content: List of markup elements for HTML/XML files
            original_file_path: Path to the original file
        """
        # PDF handling
        if file_type == "pdf" and self._can_create_pdf_split(
            original_pdf, page_numbers, chunk_indices
        ):
            # logger.debug(
            #     f"Creating PDF split from pages {[page_numbers[idx] for idx in chunk_indices if idx < len(page_numbers)]}"
            # )
            self._create_pdf_split(
                split_path, original_pdf, page_numbers, chunk_indices, content
            )

        # Spreadsheet handling
        elif file_type in ["xlsx", "csv"] and self._can_create_spreadsheet_split(
            original_spreadsheet, sheet_info, chunk_indices
        ):
            logger.debug("Creating spreadsheet split from sheet info chunks")
            self._create_spreadsheet_split(
                split_path,
                original_spreadsheet,
                sheet_info,
                chunk_indices,
                file_type,
                content,
            )

        # Markup handling
        elif file_type in ["html", "xml"] and self._can_create_markup_split(
            markup_content, chunk_indices
        ):
            # logger.debug(
            #     f"Creating markup split with {len([markup_content[idx] for idx in chunk_indices if idx < len(markup_content)])} elements"
            # )
            self._create_markup_split(
                split_path,
                original_file_path,
                markup_content,
                chunk_indices,
                file_type,
                content,
            )

        # DOCX handling
        elif file_type == "docx":
            logger.debug("Creating DOCX split")
            self._create_docx_split(split_path, content)

        # PPTX handling
        elif file_type == "pptx":
            logger.debug("Creating PPTX split")
            self._create_pptx_split(split_path, content)

        # Default: save as plain text
        else:
            # For plain text or when specialized handling fails, save as text
            logger.debug("Saving content as plain text")
            self._save_text_content(split_path, content)

    def _generate_split_path(self, original_file_path: str) -> str:
        """Generate a unique path for a split file."""
        original_filename = os.path.basename(original_file_path)
        split_id = str(uuid.uuid4())
        split_filename = f"{os.path.splitext(original_filename)[0]}_{split_id}{os.path.splitext(original_filename)[1]}"
        split_path = os.path.join(settings.TEMP_SPLITS_DIR, split_filename)
        logger.debug(
            f"Generated split path: {split_path} from original: {original_file_path}"
        )
        return split_path

    def _can_create_pdf_split(
        self, original_pdf: Path, page_numbers: List[int], chunk_indices: List[int]
    ) -> bool:
        """Check if we can create a PDF split."""
        can_create = (
            original_pdf is not None
            and page_numbers is not None
            and chunk_indices is not None
        )
        logger.debug(f"Can create PDF split: {can_create}")
        return can_create

    def _create_pdf_split(
        self,
        split_path: str,
        original_pdf: Path,
        page_numbers: List[int],
        chunk_indices: List[int],
        fallback_content: str,
    ) -> None:
        """Create a PDF split file using PyMuPDF (robust fallback to pypdf then text)."""
        # Get the relevant page numbers for this split
        relevant_pages = [
            page_numbers[idx] for idx in chunk_indices if idx < len(page_numbers)
        ]

        # Try PyMuPDF first (most robust)
        if PYMUPDF_AVAILABLE:
            try:
                self._create_pdf_split_pymupdf(split_path, original_pdf, relevant_pages)
                logger.debug(
                    f"Successfully created PDF split with PyMuPDF at: {split_path}"
                )
                return
            except Exception as e:
                logger.warning(
                    f"PyMuPDF PDF split creation failed: {str(e)}, trying pypdf fallback"
                )

        # Fallback to pypdf
        try:
            with open(original_pdf, "rb") as file:
                reader = pypdf.PdfReader(file)
                writer = pypdf.PdfWriter()

                # Add each page to the new PDF
                for page_num in relevant_pages:
                    if page_num < len(reader.pages):
                        writer.add_page(reader.pages[page_num])
                    else:
                        logger.warning(
                            f"Page {page_num} is out of range for PDF with {len(reader.pages)} pages"
                        )

                # Save the new PDF
                with open(split_path, "wb") as output_file:
                    writer.write(output_file)
                logger.debug(
                    f"Successfully created PDF split with pypdf at: {split_path}"
                )
                return
        except Exception as e:
            logger.warning(
                f"pypdf PDF split creation failed: {str(e)}, falling back to text content"
            )

        # Final fallback to text saving
        logger.info(f"Falling back to text content for PDF split at: {split_path}")
        self._save_text_content(split_path, fallback_content)

    def _create_pdf_split_pymupdf(
        self, split_path: str, original_pdf: Path, relevant_pages: List[int]
    ) -> None:
        """Create a PDF split using PyMuPDF."""
        try:
            # Open source document
            src_doc = fitz.open(str(original_pdf))

            # Create new document for the split
            dst_doc = fitz.open()

            # Copy specified pages to new document
            for page_num in relevant_pages:
                if page_num < src_doc.page_count:
                    dst_doc.insert_pdf(src_doc, from_page=page_num, to_page=page_num)
                else:
                    logger.warning(
                        f"Page {page_num} is out of range for PDF with {src_doc.page_count} pages"
                    )

            # Save the new PDF
            dst_doc.save(split_path)
            dst_doc.close()
            src_doc.close()

            logger.debug(f"PyMuPDF created PDF split: {len(relevant_pages)} pages")

        except Exception as e:
            logger.error(f"PyMuPDF PDF split creation error: {e}")
            raise

    def _can_create_spreadsheet_split(
        self,
        original_spreadsheet: Path,
        sheet_info: List[Tuple[str, int, int]],
        chunk_indices: List[int],
    ) -> bool:
        """Check if we can create a spreadsheet split."""
        return (
            original_spreadsheet is not None
            and sheet_info is not None
            and chunk_indices is not None
        )

    def _create_spreadsheet_split(
        self,
        split_path: str,
        original_spreadsheet: Path,
        sheet_info: List[Tuple[str, int, int]],
        chunk_indices: List[int],
        file_type: str,
        fallback_content: str,
    ) -> None:
        """Create a spreadsheet split file (CSV or Excel)."""
        try:
            if file_type == "csv":
                self._create_csv_split(
                    split_path, original_spreadsheet, sheet_info, chunk_indices
                )
            else:  # Excel
                self._create_excel_split(
                    split_path, original_spreadsheet, sheet_info, chunk_indices
                )
        except Exception as e:
            logger.error(f"Error creating spreadsheet split: {e}")
            # Fall back to text saving if spreadsheet creation fails
            self._save_text_content(split_path, fallback_content)

    def _create_csv_split(
        self,
        split_path: str,
        original_spreadsheet: Path,
        sheet_info: List[Tuple[str, int, int]],
        chunk_indices: List[int],
    ) -> None:
        """Create a CSV split file with encoding detection."""
        # Detect encoding for the original CSV
        encoding = self.detect_file_encoding(original_spreadsheet)

        try:
            # For CSV, filter the rows and save as a new CSV
            df = pd.read_csv(original_spreadsheet, encoding=encoding, low_memory=False)
        except Exception as e:
            logger.warning(
                f"Failed to read CSV with detected encoding {encoding}, trying fallbacks: {str(e)}"
            )
            # Try common encodings as fallback
            for fallback_encoding in ["utf-8", "latin-1", "cp1252"]:
                try:
                    df = pd.read_csv(
                        original_spreadsheet,
                        encoding=fallback_encoding,
                        low_memory=False,
                    )
                    logger.info(
                        f"Successfully read original CSV with fallback encoding: {fallback_encoding}"
                    )
                    encoding = fallback_encoding
                    break
                except Exception as e:
                    logger.warning(
                        f"Failed to read CSV with fallback encoding {fallback_encoding}: {str(e)}"
                    )
                    continue
            else:
                # Last resort: let pandas auto-detect
                df = pd.read_csv(original_spreadsheet, encoding=None, low_memory=False)
                encoding = "utf-8"  # Default for output

        # Get the relevant row ranges for this split
        relevant_ranges = [
            sheet_info[idx] for idx in chunk_indices if idx < len(sheet_info)
        ]

        rows = []
        for _, start_row, end_row in relevant_ranges:
            rows.extend(list(range(start_row, end_row)))

        # Create a new DataFrame with only the selected rows
        df_split = df.iloc[rows]

        # Save with UTF-8 encoding to ensure compatibility
        df_split.to_csv(split_path, index=False, encoding="utf-8")

    def _create_excel_split(
        self,
        split_path: str,
        original_spreadsheet: Path,
        sheet_info: List[Tuple[str, int, int]],
        chunk_indices: List[int],
    ) -> None:
        """Create an Excel split file."""
        # For Excel, create a new Excel file with selected sheets and rows
        sheets = pd.read_excel(original_spreadsheet, sheet_name=None)
        writer = pd.ExcelWriter(split_path, engine="openpyxl")

        # Process each sheet
        processed_sheets = set()
        for idx in chunk_indices:
            if idx < len(sheet_info):
                sheet_name, start_row, end_row = sheet_info[idx]
                if sheet_name in sheets and sheet_name not in processed_sheets:
                    df = sheets[sheet_name]
                    # If this is the first chunk for this sheet, create a new sheet
                    df_split = df.iloc[start_row:end_row]
                    df_split.to_excel(writer, sheet_name=sheet_name, index=False)
                    processed_sheets.add(sheet_name)

        writer.close()

    def _can_create_markup_split(
        self, markup_content: List[str], chunk_indices: List[int]
    ) -> bool:
        """Check if we can create a markup split."""
        return markup_content is not None and chunk_indices is not None

    def _create_markup_split(
        self,
        split_path: str,
        original_file_path: str,
        markup_content: List[str],
        chunk_indices: List[int],
        file_type: str,
        fallback_content: str,
    ) -> None:
        """Create a markup split file (HTML or XML)."""
        try:
            # Get the relevant markup elements
            original_filename = os.path.basename(original_file_path)
            relevant_markup = [
                markup_content[idx]
                for idx in chunk_indices
                if idx < len(markup_content)
            ]

            if file_type == "html":
                self._create_html_split(split_path, original_filename, relevant_markup)
            else:  # XML
                self._create_xml_split(split_path, relevant_markup)
        except Exception as e:
            logger.error(f"Error creating markup split: {e}")
            # Fall back to text saving if markup creation fails
            self._save_text_content(split_path, fallback_content)

    def _create_html_split(
        self, split_path: str, original_filename: str, relevant_markup: List[str]
    ) -> None:
        """Create an HTML split file."""
        # Create a basic HTML document structure
        html_content = f"""<!DOCTYPE html>
<html>
<head>
    <title>Split from {original_filename}</title>
</head>
<body>
    {"".join(relevant_markup)}
</body>
</html>"""
        with open(split_path, "w", encoding="utf-8") as f:
            f.write(html_content)

    def _create_xml_split(self, split_path: str, relevant_markup: List[str]) -> None:
        """Create an XML split file."""
        # Create a simple XML wrapper around the elements
        xml_content = f"""<?xml version="1.0" encoding="UTF-8"?>
<root>
    {"".join(relevant_markup)}
</root>"""
        with open(split_path, "w", encoding="utf-8") as f:
            f.write(xml_content)

    def _create_docx_split(self, split_path: str, content: str) -> None:
        """Create a DOCX split file."""
        try:
            # Create a new DOCX with the content
            doc = Document()
            for paragraph in content.split("\n\n"):
                if paragraph.strip():
                    doc.add_paragraph(paragraph)
            doc.save(split_path)
        except Exception as e:
            logger.error(f"Error creating DOCX split: {e}")
            # Fall back to text saving if DOCX creation fails
            self._save_text_content(split_path, content)

    def _create_pptx_split(self, split_path: str, content: str) -> None:
        """Create a PPTX split file."""
        try:
            # Create a new PowerPoint presentation with the content
            prs = Presentation()

            # Split content by slides (assuming content contains slide markers)
            slides_content = content.split("Slide ")

            for slide_content in slides_content:
                if slide_content.strip():
                    # Create a new slide
                    slide_layout = prs.slide_layouts[1]  # Use title and content layout
                    slide = prs.slides.add_slide(slide_layout)

                    # Parse slide content to extract title and body
                    lines = slide_content.strip().split("\n")
                    if lines:
                        # First line becomes the title (remove slide number if present)
                        title_text = lines[0]
                        if ":" in title_text:
                            title_text = title_text.split(":", 1)[1].strip()

                        # Set slide title
                        if slide.shapes.title:
                            slide.shapes.title.text = title_text or "Slide Content"

                        # Add remaining content to the body
                        if len(lines) > 1:
                            body_text = "\n".join(lines[1:]).strip()
                            if body_text and hasattr(slide.placeholders[1], "text"):
                                slide.placeholders[1].text = body_text

            # If no slides were created, create a single slide with all content
            if len(prs.slides) == 0:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = "Content"
                if hasattr(slide.placeholders[1], "text"):
                    slide.placeholders[1].text = content

            prs.save(split_path)

        except Exception as e:
            logger.error(f"Error creating PPTX split: {e}")
            # Fall back to text saving if PPTX creation fails
            self._save_text_content(split_path, content)

    def _save_text_content(self, split_path: str, content: str) -> None:
        """Save content as a plain text file."""
        with open(split_path, "w", encoding="utf-8") as f:
            f.write(content)
        logger.debug(f"Saved text content to {split_path} ({len(content)} characters)")

    def estimate_tokens_for_input(self, input_data: str) -> int:
        """
        Estimate the number of tokens in an input string.

        Args:
            input_data: The input string to count tokens for

        Returns:
            Estimated number of tokens
        """
        logger.debug(f"Estimating tokens for input string of length {len(input_data)}")
        try:
            token_count = self.count_tokens(input_data)
            logger.debug(f"Calculated token count: {token_count}")
            return token_count
        except Exception as e:
            logger.error(f"Error estimating tokens for input: {str(e)}")
            # Fall back to character-based estimation
            # Rough estimate: 1 token per 4 characters
            token_count = len(input_data) // 4
            logger.warning(
                f"Falling back to character-based estimation: {token_count} tokens"
            )
            return token_count
